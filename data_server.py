import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import hashlib
import json
import os

from mcp.server.fastmcp import FastMCP
from mcp.server.fastmcp.prompts import base
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder
from sklearn.metrics import accuracy_score, root_mean_squared_error
from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from dotenv import load_dotenv
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI

load_dotenv()
os.getenv("OPENAI_API_KEY")

mcp = FastMCP("DataAnalysis")

ppt_folder = "./pptx"
index_folder = "faiss_ppt_index"
meta_path = os.path.join(index_folder, "index_meta.json")

def get_ppt_folder_hash(folder_path: str) -> str:
    """
    현재 pptx 폴더 내 파일명 + 수정시간 기준으로 해시를 계산합니다.
    """
    files = sorted([
        (f, os.path.getmtime(os.path.join(folder_path, f)))
        for f in os.listdir(folder_path)
        if f.endswith(".pptx")
    ])
    hash_input = json.dumps(files).encode()
    return hashlib.md5(hash_input).hexdigest()


@mcp.tool()
def search_ppt(query: str) -> str:
    """
    Search and summarize the content of PowerPoint (.pptx) files related to a given question.
    
    This tool loads PowerPoint slides stored in the ./pptx directory, indexes the text using a vector database (FAISS),
    and uses a language model to answer the user's question based on the content of the slides.
    
    Args:
        query (str): A natural language question to ask about the slide content.

    Returns:
        str: The LLM-generated answer based on slide content.
    """

    ppt_hash = get_ppt_folder_hash(ppt_folder)
    need_update = True

    if os.path.exists(index_folder) and os.path.exists(meta_path):
        try:
            with open(meta_path) as f:
                saved_hash = json.load(f).get("ppt_hash")
                if saved_hash == ppt_hash:
                    need_update = False
        except Exception:
            pass

    if need_update:
        all_text = []
        for fname in os.listdir(ppt_folder):
            if fname.endswith(".pptx"):
                prs = Presentation(os.path.join(ppt_folder, fname))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            all_text.append(shape.text)

        text = "\n".join(all_text)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        docs = text_splitter.create_documents([text])

        vectordb = FAISS.from_documents(docs, OpenAIEmbeddings())
        vectordb.save_local(index_folder)

        os.makedirs(index_folder, exist_ok=True)
        with open(meta_path, "w") as f:
            json.dump({"ppt_hash": ppt_hash}, f)
    else:
        vectordb = FAISS.load_local(index_folder, OpenAIEmbeddings())

    retriever = vectordb.as_retriever()
    qa = RetrievalQA.from_chain_type(llm=ChatOpenAI(model="gpt-4o"), retriever=retriever)
    result = qa.run(query)
    return result

@mcp.tool()
def describe_column(csv_path: str, column: str) -> dict:
    """
    Get summary statistics (count, mean, std, min, max, etc.) for a specific column in a CSV file.

    Args:
        csv_path (str): The file path to the CSV file.
        column (str): The name of the column to compute statistics for.

    Returns:
        dict: A dictionary containing summary statistics of the specified column.
    """
    df = pd.read_csv(csv_path)
    if column not in df.columns:
        raise ValueError(f"Column '{column}' not found in CSV.")
    return df[column].describe().to_dict()


@mcp.tool()
def plot_histogram(csv_path: str, column: str, bins: int = 10) -> str:
    """
    Generate and save a density histogram for a specific column in a CSV file.

    Args:
        csv_path (str): The file path to the CSV file.
        column (str): The name of the column to visualize.
        bins (int, optional): Number of histogram bins. Defaults to 10.

    Returns:
        str: The file path to the saved density histogram image.
    """
    df = pd.read_csv(csv_path)
    if column not in df.columns:
        raise ValueError(f"Column '{column}' not found in CSV.")

    plt.figure(figsize=(8, 6))
    sns.histplot(
        df[column].dropna(),
        bins=bins,
        kde=True,
        stat="density",
        edgecolor="black",
        alpha=0.6,
    )
    plt.xlabel(column)
    plt.ylabel("Density")
    plt.title(f"Density Histogram of {column}")

    output_path = f"{column}_density_hist.png"
    plt.savefig(output_path)
    plt.close()

    return output_path


@mcp.tool()
def model(csv_path: str, x_columns: list, y_column: str) -> dict:
    """
    Automatically train a model (classification or regression) based on the target column type.

    Args:
        csv_path: Path to CSV file.
        x_columns: List of feature column names.
        y_column: Target column name.

    Returns:
        Dictionary with model type, performance metric, and score.
    """
    df = pd.read_csv(csv_path)

    for col in x_columns + [y_column]:
        if col not in df.columns:
            raise ValueError(f"Column '{col}' not found in CSV.")

    X = df[x_columns]
    y = df[y_column]

    for col in X.select_dtypes(include=["object"]).columns:
        X[col] = LabelEncoder().fit_transform(X[col])

    is_classification = y.dtype == "object" or len(y.unique()) <= 10

    if is_classification:
        y = LabelEncoder().fit_transform(y)
        model = RandomForestClassifier()
        metric_name = "accuracy"
    else:
        model = RandomForestRegressor()
        metric_name = "rmse"

    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=42
    )

    model.fit(X_train, y_train)
    y_pred = model.predict(X_test)

    if is_classification:
        score = accuracy_score(y_test, y_pred)
        model_type = "classification"
    else:
        score = root_mean_squared_error(y_test, y_pred, squared=False)
        model_type = "regression"

    return {"model_type": model_type, "metric": metric_name, "score": score}


@mcp.prompt()
def default_prompt(message: str) -> list[base.Message]:
    return [
        base.AssistantMessage(
            "You are a helpful data analysis assistant. \n"
            "Please clearly organize and return the results of the tool calling and the data analysis."
        ),
        base.UserMessage(message),
    ]


if __name__ == "__main__":
    mcp.run(transport="stdio")